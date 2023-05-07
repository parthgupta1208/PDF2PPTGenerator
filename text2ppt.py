from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
#Create a new PowerPoint presentation

def presentate(defined_list):
    prs = Presentation()

    def add_slide(prs, layout, title, subtitle):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text=subtitle
        font = slide.shapes.title.text_frame.paragraphs[0].font
        font.name = 'Arial'
        font.size = Pt(30)
        font.bold = True
        font.italic = False
        for x in  slide.placeholders[1].text_frame.paragraphs:
            font1= x.font
            font1.name = 'Arial'
            font1.size = Pt(16)
            font1.bold = False
            font1.italic = False
        return slide


    def add_slide_img(prs, layout, img_path):
        slide = prs.slides.add_slide(layout)
        img_path =  ""+img_path
        left =  Inches(1.10)
        top = Inches(0.7)
        width = Inches(8)
        height = Inches(6)
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
        
    title_slide_layout = prs.slide_layouts[1]
    title_slide_layimg=prs.slide_layouts[6]
    for i in range (0,len(defined_list)):     
        slide = add_slide(prs, title_slide_layout, defined_list[i]["Topic"],"\n".join(defined_list[i]["Summary"]))

    # Save the presentation
    prs.save("PPT.pptx")